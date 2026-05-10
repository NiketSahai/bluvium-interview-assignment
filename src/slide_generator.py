"""
Slide Deck Generator module.

Produces a python-pptx slide deck embedding charts and key findings
from the Transcript Intelligence pipeline.
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Color scheme matching the visualization engine
COLORS = {
    'primary_blue': RGBColor(0x1F, 0x77, 0xB4),
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),
    'orange': RGBColor(0xFF, 0x7F, 0x0E),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'medium_gray': RGBColor(0x66, 0x66, 0x66),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_NAME = 'Calibri'


def _add_title_and_subtitle(slide, title_text: str, subtitle_text: str = ''):
    """Add title and optional subtitle to a slide using placeholders."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT_NAME
                run.font.color.rgb = COLORS['dark_blue']

    if subtitle_text:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = subtitle_text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = FONT_NAME
                        run.font.color.rgb = COLORS['medium_gray']
                break


def _add_content_slide(prs, title: str, bullets: list[str],
                       stakeholder: str = '') -> None:
    """Add a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(28)
                run.font.color.rgb = COLORS['dark_blue']

    # Set content bullets
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS['dark_gray']
                p.space_after = Pt(6)
            break

    # Add stakeholder attribution at bottom if provided
    if stakeholder:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"For: {stakeholder}"
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLORS['orange']


def _add_chart_slide(prs, title: str, figure_path: str,
                     stakeholder: str = '') -> None:
    """Add a slide with a chart image taking up most of the slide area."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title at top
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # Add chart image or placeholder
    if figure_path and os.path.exists(figure_path):
        slide.shapes.add_picture(
            figure_path,
            Inches(0.5), Inches(1.0),
            Inches(9.0), Inches(5.8)
        )
    else:
        # Add placeholder text if image is missing
        txBox = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.0), Inches(7.0), Inches(1.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Chart not available: {os.path.basename(figure_path) if figure_path else 'no path provided'}]"
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = COLORS['medium_gray']
        p.alignment = PP_ALIGN.CENTER

    # Add stakeholder attribution at bottom if provided
    if stakeholder:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.9), Inches(9.0), Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"For: {stakeholder}"
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLORS['orange']


def _get_finding(findings: dict, key: str, default: str = '') -> str:
    """Safely get a finding value with a default fallback."""
    return findings.get(key, default) or default


def _split_into_bullets(text: str, max_bullets: int = 6) -> list[str]:
    """Split a findings text into bullet points."""
    if not text:
        return ['No data available']

    # If text already has bullet-like structure (newlines), split on them
    lines = [line.strip() for line in text.split('\n') if line.strip()]

    if len(lines) > 1:
        return lines[:max_bullets]

    # Otherwise split on sentences/periods
    sentences = [s.strip() for s in text.split('.') if s.strip()]
    bullets = [s + '.' if not s.endswith('.') else s for s in sentences]
    return bullets[:max_bullets]


def generate_slide_deck(
    figures: dict[str, str],
    findings: dict[str, str],
    output_path: str = 'output/transcript_intelligence_deck.pptx'
) -> None:
    """
    Generate presentation slide deck.

    Args:
        figures: Dict mapping chart_name -> figure path (PNG files in output/figures/)
        findings: Dict mapping section_name -> findings text (key insights to display)
        output_path: Path to save the .pptx file
    """
    # Ensure output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir:
        Path(output_dir).mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Set slide dimensions to standard 10" x 7.5"
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    _add_title_and_subtitle(
        slide,
        'Transcript Intelligence',
        'Insights from 100 AegisCloud Meeting Transcripts'
    )

    # --- Slide 2: Executive Summary ---
    exec_summary = _get_finding(
        findings, 'executive_summary',
        'Analysis of 100 meeting transcripts reveals key patterns in customer sentiment, '
        'feature requests, incident response, and churn risk across AegisCloud stakeholders.'
    )
    _add_content_slide(
        prs,
        'Executive Summary',
        _split_into_bullets(exec_summary),
        stakeholder='Product & Engineering Leadership'
    )

    # --- Slide 3: Methodology Overview ---
    methodology = _get_finding(
        findings, 'methodology',
        'Pipeline uses VADER sentiment analysis, rule-based topic categorization with TF-IDF validation, '
        'and keyword-based insight extraction.\n'
        'All processing runs locally without external API dependencies.\n'
        'Meetings classified by call type: internal, external, support.\n'
        'Key moments extracted from pre-labeled summary data.'
    )
    _add_content_slide(
        prs,
        'Methodology Overview',
        _split_into_bullets(methodology)
    )

    # --- Slide 4: Dataset Overview ---
    dataset_overview = _get_finding(
        findings, 'dataset_overview',
        '100 meeting transcripts analyzed.\n'
        'Call types: Internal, External, Support.\n'
        'Date range spans multiple months of AegisCloud operations.\n'
        'Data includes transcripts, summaries, key moments, and action items.'
    )
    _add_content_slide(
        prs,
        'Dataset Overview',
        _split_into_bullets(dataset_overview)
    )

    # --- Slides 5-6: Topic Categorization ---
    topic_findings = _get_finding(
        findings, 'topic_findings',
        'Topics categorized across Product areas (Detect, Comply, Identity) '
        'and operational themes (incident response, sprint planning, customer success).'
    )
    _add_content_slide(
        prs,
        'Topic Categorization: Key Findings',
        _split_into_bullets(topic_findings),
        stakeholder='Product Managers'
    )

    _add_chart_slide(
        prs,
        'Topic Distribution',
        figures.get('topic_distribution', ''),
        stakeholder='Product Managers'
    )

    # --- Slides 7-8: Sentiment Analysis ---
    sentiment_findings = _get_finding(
        findings, 'sentiment_findings',
        'Sentiment analysis reveals patterns across call types.\n'
        'VADER scores correlated with pre-existing sentiment labels.\n'
        'Time-series trends show sentiment progression over the analysis period.'
    )
    _add_content_slide(
        prs,
        'Sentiment Analysis: Key Findings',
        _split_into_bullets(sentiment_findings),
        stakeholder='Support Leaders'
    )

    _add_chart_slide(
        prs,
        'Sentiment Over Time by Call Type',
        figures.get('sentiment_time_series', ''),
        stakeholder='Support Leaders'
    )

    # --- Slides 9-10: Churn Risk & Renewal Intelligence ---
    churn_findings = _get_finding(
        findings, 'churn_findings',
        'Churn risk assessed for external accounts based on signal frequency, '
        'negative sentiment, and competitive mentions.\n'
        'At-risk accounts ranked by composite risk score (0-10 scale).'
    )
    _add_content_slide(
        prs,
        'Churn Risk & Renewal Intelligence',
        _split_into_bullets(churn_findings),
        stakeholder='Sales Managers'
    )

    _add_chart_slide(
        prs,
        'Churn Risk Ranking: Top At-Risk Accounts',
        figures.get('churn_risk_ranking', ''),
        stakeholder='Sales Managers'
    )

    # --- Slides 11-12: Feature Gap Analysis ---
    feature_gap_findings = _get_finding(
        findings, 'feature_gap_findings',
        'Feature gaps extracted from all call types and grouped by product area.\n'
        'Gaps ranked by frequency of mention across distinct meetings.\n'
        'Stakeholder attribution captured for each gap.'
    )
    _add_content_slide(
        prs,
        'Feature Gap Analysis',
        _split_into_bullets(feature_gap_findings),
        stakeholder='Product Managers'
    )

    _add_chart_slide(
        prs,
        'Feature Gap Frequency by Product Area',
        figures.get('feature_gap_frequency', ''),
        stakeholder='Product Managers'
    )

    # --- Slides 13-14: Incident Patterns & Team Health ---
    incident_findings = _get_finding(
        findings, 'incident_findings',
        'Incident patterns identified from internal meetings.\n'
        'Recurring issues flagged as systemic (3+ meeting mentions).\n'
        'Team health indicators extracted: timeline risk, resource constraints, process gaps.'
    )
    _add_content_slide(
        prs,
        'Incident Patterns & Team Health',
        _split_into_bullets(incident_findings),
        stakeholder='Engineering Leads'
    )

    _add_chart_slide(
        prs,
        'Incident Timeline',
        figures.get('incident_timeline', ''),
        stakeholder='Engineering Leads'
    )

    # --- Slides 15-16: Support Issue Analysis ---
    support_findings = _get_finding(
        findings, 'support_findings',
        'Support issues categorized by type: billing, technical, provisioning, compliance.\n'
        'High-frustration cases identified via negative sentiment threshold.\n'
        'Resolution patterns detected through positive_pivot key moments.'
    )
    _add_content_slide(
        prs,
        'Support Issue Analysis',
        _split_into_bullets(support_findings),
        stakeholder='Support Leaders'
    )

    _add_chart_slide(
        prs,
        'Support Issue Categories',
        figures.get('support_categories', ''),
        stakeholder='Support Leaders'
    )

    # --- Slide 17: Recommendations ---
    recommendations = _get_finding(
        findings, 'recommendations',
        'Prioritize retention outreach for high-risk accounts.\n'
        'Address top feature gaps in product roadmap planning.\n'
        'Investigate systemic incidents for root cause resolution.\n'
        'Optimize support workflows for high-frustration categories.\n'
        'Continue monitoring sentiment trends across all call types.'
    )
    _add_content_slide(
        prs,
        'Recommendations',
        _split_into_bullets(recommendations),
        stakeholder='Product & Engineering Leadership'
    )

    # --- Slide 18: Appendix / Methodology Details ---
    appendix_bullets = [
        'Sentiment: VADER compound scores normalized to 1-5 scale',
        'Topics: Rule-based heuristics validated by TF-IDF clustering',
        'Churn Risk: Composite score from signals, sentiment, competitive mentions',
        'Feature Gaps: Keyword-based product area assignment (Detect, Comply, Identity, Platform)',
        'Incidents: Systemic threshold at 3+ meeting mentions',
        'Support: Category assignment via title and summary keyword matching',
        'All processing local — no external API dependencies',
    ]
    _add_content_slide(
        prs,
        'Appendix: Methodology Details',
        appendix_bullets
    )

    # Save the presentation
    prs.save(output_path)
